#!/usr/bin/env python3
"""
Analyze the Hard Rock Digital PowerPoint template
Extract layouts, colors, fonts, and design elements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

template_path = "/home/user/test/HRD_PowerPoint-Template_v1.pptx"

print("="*80)
print("ANALYZING HARD ROCK DIGITAL POWERPOINT TEMPLATE")
print("="*80)

prs = Presentation(template_path)

print(f"\n📄 Presentation Properties:")
print(f"   Slide Width: {prs.slide_width / 914400:.2f} inches")
print(f"   Slide Height: {prs.slide_height / 914400:.2f} inches")
print(f"   Number of slides: {len(prs.slides)}")
print(f"   Number of slide layouts: {len(prs.slide_layouts)}")

print(f"\n📐 Available Slide Layouts:")
for idx, layout in enumerate(prs.slide_layouts):
    print(f"   [{idx}] {layout.name}")
    print(f"       Placeholders: {len(layout.placeholders)}")
    for pidx, placeholder in enumerate(layout.placeholders):
        try:
            print(f"         - [{placeholder.placeholder_format.idx}] {placeholder.name} ({placeholder.placeholder_format.type})")
        except:
            print(f"         - Placeholder {pidx}")

print(f"\n🎨 Analyzing Slides in Template:")
for slide_idx, slide in enumerate(prs.slides):
    print(f"\n   Slide {slide_idx + 1}:")

    # Check background
    try:
        if slide.background.fill.type:
            print(f"      Background: {slide.background.fill.type}")
            if hasattr(slide.background.fill, 'fore_color'):
                try:
                    rgb = slide.background.fill.fore_color.rgb
                    print(f"      Background Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]}) / #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    print(f"      Background Color: (theme-based)")
    except:
        print(f"      Background: Standard")

    # Analyze shapes
    print(f"      Shapes: {len(slide.shapes)}")
    for shape_idx, shape in enumerate(slide.shapes):
        shape_type = type(shape).__name__
        print(f"         [{shape_idx}] {shape_type}: {shape.name if hasattr(shape, 'name') else 'unnamed'}")

        # Check if it has text
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            text = shape.text_frame.text[:50] if shape.text_frame.text else "(empty)"
            print(f"             Text: {text}")

            # Check font properties
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        print(f"             Font: {run.font.name if run.font.name else 'default'}")
                        print(f"             Size: {run.font.size.pt if run.font.size else 'default'} pt")
                        try:
                            if run.font.color.type:
                                rgb = run.font.color.rgb
                                print(f"             Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]}) / #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                        except:
                            print(f"             Color: theme-based")
                        break
                break

        # Check fill color
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type:
                    print(f"             Fill Type: {shape.fill.type}")
                    if hasattr(shape.fill, 'fore_color'):
                        try:
                            rgb = shape.fill.fore_color.rgb
                            print(f"             Fill Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]}) / #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                        except:
                            pass
            except:
                pass

print("\n" + "="*80)
print("ANALYSIS COMPLETE")
print("="*80)
