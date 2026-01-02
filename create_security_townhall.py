#!/usr/bin/env python3
"""
Create State of Security 2026 Townhall Presentation
Following Hard Rock Digital Brand Guidelines
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors from guidelines
VIOLET_1 = RGBColor(106, 70, 243)  # Primary brand color
VIOLET_2 = RGBColor(195, 37, 180)   # Secondary brand color
BLUE_1 = RGBColor(15, 197, 222)
BLUE_2 = RGBColor(16, 4, 88)
BLUE_3 = RGBColor(63, 133, 238)
WHITE = RGBColor(248, 248, 250)
BLACK = RGBColor(26, 24, 27)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with brand styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = VIOLET_1

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = VIOLET_1
    title_shape.line.fill.background()

    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_before = Pt(12)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = VIOLET_1
    title_shape.line.fill.background()

    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.5)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(10)

    return slide

# Slide 1: Title Slide
add_title_slide(prs,
    "State of Security 2026",
    "Building Tomorrow's Security | Hard Rock Digital Security Townhall"
)

# Slide 2: Agenda
add_content_slide(prs, "Agenda", [
    "Year in Review: 2025 Achievements",
    "Deep Dive: AI-Powered Security Operations",
    "Security Operations & Detection Engineering",
    "Governance & Maturity",
    "Current State: Where We Stand Today",
    "Future State: Zero Trust Vision",
    "2026 Priorities & 3-Year Roadmap",
    "Our Philosophy: Partnership & Risk-Based Security",
    "Q&A"
])

# Slide 3: Our Security Philosophy
add_content_slide(prs, "Our Security Philosophy", [
    "🤝 Risk-Based Approach",
    "   • We inform, business decides",
    "   • Security highlights risks, leadership accepts/mitigates",
    "",
    "🚀 Business Enablement",
    "   • Go as fast as you can possibly go",
    "   • We remove blockers, not create them",
    "",
    "💡 FAIL = First Attempt In Learning",
    "   • Fail early, fail often, fail fast",
    "   • Learning mindset over perfection",
    "",
    "🤝 Partnership, Not Gatekeeping",
    "   • We don't approve/deny—we collaborate"
])

# Slide 4: Executive Summary: 2025 by the Numbers
add_content_slide(prs, "2025 By The Numbers", [
    "📊 Key Metrics:",
    "   • 500 analytics rules deployed (at platform capacity)",
    "   • 79.7% alert closure rate maintained",
    "   • 49+ vendor security reviews completed",
    "   • ~560 employees trained on security awareness",
    "   • 33% cost reduction in Sentinel logging",
    "",
    "👥 Team Growth:",
    "   • Director of SRM hired",
    "   • 3 new SOC analysts",
    "   • 2 new SRM analysts",
    "   • 1 Security Engineer (Jan 2026)",
    "",
    "🌎 Market Launches:",
    "   • Colorado & Michigan successfully launched",
    "   • Ontario launch prep underway (Q1 2026)"
])

# Slide 5: AI & MCP: Leading the Industry (Part 1)
add_content_slide(prs, "🤖 AI & MCP: Industry Pioneers", [
    "What is Model Context Protocol (MCP)?",
    "   • Standardized way for AI to interact with external systems",
    "   • Real-time access to security tools and data",
    "   • Natural language interface for complex operations",
    "",
    "Why It Matters:",
    "   • We're among the FIRST security teams deploying production MCP",
    "   • Transforms how analysts interact with security tools",
    "   • Reduces manual work, accelerates investigation",
    "   • Enables AI-assisted decision making",
    "",
    "Our Position:",
    "   • Industry leadership in AI-powered security operations",
    "   • Innovation differentiator for Hard Rock Digital",
    "   • Competitive advantage in threat detection & response"
])

# Slide 6: AI & MCP: Our Implementations
add_two_column_slide(prs, "🤖 AI & MCP: Our Implementations",
    [
        "TheHive MCP Server:",
        "• Natural language case management",
        "• AI-assisted investigation workflow",
        "• Real-time alert context generation",
        "• Automated triage recommendations",
        "",
        "Vanta MCP Server:",
        "• Programmatic compliance access",
        "• Automated security posture reporting",
        "• Real-time compliance status queries",
        "• Policy violation detection",
        "",
        "AI SOC Level 1 'Analyst':",
        "• Initial alert triage automation",
        "• Pattern recognition across alerts",
        "• Reduces analyst burnout on repetitive tasks"
    ],
    [
        "Current AI Capabilities:",
        "• Claude AI integration for analysis",
        "• AI-powered log analysis",
        "• RSS cyber threat summarization",
        "• Policy development assistance",
        "• Security documentation generation",
        "",
        "2026 AI Roadmap:",
        "• Expand MCP to additional platforms",
        "• AI-powered playbook automation",
        "• Custom threat detection models",
        "• Enhanced behavioral analytics",
        "• BurpGPT for security testing",
        "",
        "Governance:",
        "• Hallucination testing completed",
        "• Least privilege OAuth credentials",
        "• Security evaluation framework"
    ]
)

# Slide 7: Security Operations Excellence
add_content_slide(prs, "Security Operations Excellence", [
    "🔒 24/7 Monitoring - Continent8 MSOC Partnership:",
    "   • 13 SOC analyst accounts onboarded",
    "   • API integration for automated alert forwarding",
    "   • n8n workflow automation (6pm-2am EST coverage)",
    "",
    "📈 Detection Engineering:",
    "   • 500 analytics rules at platform capacity",
    "   • 79.7% alert closure rate",
    "   • Major tuning: AWS false positives reduced from 647 to 4 events",
    "   • New capabilities: DPRK email detection, process hollowing, OAuth monitoring",
    "",
    "🎯 Incident Response Wins:",
    "   • GraphQL scraping campaign: Multiple attacks blocked, zero customer impact",
    "   • Credential stuffing (Dec 2025): 857 users rate-limited, no compromises",
    "   • Super Bowl 2025 prep: Cribl+AKS scalability tested for 10-20x load"
])

# Slide 8: Governance & Maturity Milestones
add_content_slide(prs, "Governance & Maturity Milestones", [
    "📋 ISMS Policies Approved (December 2025):",
    "   • Security & Privacy Steering Committee re-established",
    "   • Top Level Security Policy, ISMS Roles & Responsibilities",
    "   • Patch Management SLAs, Backup Policy, Vulnerability Scanning",
    "",
    "🏆 ISO 22301 BCMS Framework:",
    "   • Comprehensive Business Continuity Management System",
    "   • 10 BCMS documents completed",
    "   • RTO/RPO targets established (Tier 1: <1hr, Tier 2: <4hrs, Tier 3: <24hrs)",
    "   • Gaming Operations, Payment Processing, Cybersecurity plans",
    "",
    "✅ ISO 27001 Gap Analysis:",
    "   • 96 controls mapped across 36 systems",
    "   • RACI matrix development",
    "   • Risk treatment roadmap established",
    "   • Foundation for ISO 27001 certification (aspirational)"
])

# Slide 9: Current State: Balanced View
add_two_column_slide(prs, "Current State: Where We Stand",
    [
        "✅ Strengths:",
        "• Governance maturity (ISMS approved)",
        "• AI/MCP industry leadership",
        "• 24/7 monitoring with MSOC",
        "• Strong team growth (6 new hires)",
        "• Cost optimization (33% Sentinel savings)",
        "• Business enablement (CO, MI launches)",
        "• 79.7% alert closure rate",
        "• Comprehensive training (~560 employees)",
        "",
        "🎯 Mature Capabilities:",
        "• Detection engineering",
        "• Incident response",
        "• Vendor risk management",
        "• Business continuity planning",
        "• Compliance frameworks"
    ],
    [
        "🎯 Opportunities (2026 Focus):",
        "• Sentinel capacity (at 500 rule limit)",
        "• Platform integrations:",
        "   - Teleskope DLP (Q1 2026)",
        "   - Sublime Security email (Q1 2026)",
        "   - Flare.io dark web (Q1 2026)",
        "• Identity governance (Entra ID rollout)",
        "• Penetration testing cadence",
        "• Ontario launch security readiness",
        "",
        "💡 Not Weaknesses—Growth Areas:",
        "• Every gap has a funded plan",
        "• Proactive identification shows maturity",
        "• Systematic approach to improvement",
        "• Risk-based prioritization"
    ]
)

# Slide 10: Future State Vision: Zero Trust
add_content_slide(prs, "Future State Vision: Zero Trust", [
    "🎯 Zero Trust Principles:",
    "   • Never trust, always verify",
    "   • Identity-based access (not perimeter-based)",
    "   • Least privilege by default",
    "   • Continuous authentication & authorization",
    "",
    "✅ Current Progress:",
    "   • Cloudflare Zero Trust implementation (identity-based access)",
    "   • Entra Conditional Access policies (12-hour session limits)",
    "   • Admin account concept across teams",
    "   • Island Enterprise Browser (Phase 1 deployed)",
    "",
    "🚀 The Journey Ahead:",
    "   • Complete Island Browser enterprise rollout",
    "   • Expand identity-based controls",
    "   • Implement continuous authentication",
    "   • Enhance micro-segmentation",
    "   • Identity governance with Entra ID"
])

# Slide 11: 2026 Priorities (Q1-Q2)
add_content_slide(prs, "2026 Priorities: Q1-Q2", [
    "🌎 Ontario Market Launch Security Readiness",
    "   • Complete security assessments & regulatory documentation",
    "   • Ensure compliance with Ontario gaming regulations",
    "",
    "🎯 ISO 22301 BCMS Execution",
    "   • Execute recovery strategies, conduct tabletop exercises",
    "   • Validate RTO/RPO targets through live testing",
    "   • Pursue ISO 22301 certification (competitive differentiator)",
    "",
    "📊 Sentinel Capacity Resolution",
    "   • Address 500 analytics rule limit blocking new detections",
    "   • Evaluate Azure Data Explorer for cost-effective expansion",
    "",
    "🔒 Penetration Testing Program",
    "   • TrustedSec LLM assessment (January 20, 2026)",
    "   • Establish recurring penetration test cadence",
    "",
    "🚀 New Platform Integrations",
    "   • Teleskope DLP, Sublime Security email, Flare.io dark web"
])

# Slide 12: 3-Year Vision (2026-2028)
add_content_slide(prs, "3-Year Vision: 2026-2028", [
    "🔑 Identity Security Focus:",
    "   • Comprehensive identity governance (Entra ID)",
    "   • Quarterly User Access Reviews across 20+ systems",
    "   • Role-based access control (RBAC) maturity",
    "   • Automated access certification workflows",
    "",
    "📧 Email Security Maturity:",
    "   • Sublime Security full deployment",
    "   • Advanced phishing protection",
    "   • Business email compromise prevention",
    "",
    "🌍 International Expansion Readiness:",
    "   • Adapt to shifting business priorities",
    "   • Security frameworks for new jurisdictions beyond Ontario",
    "",
    "🏆 ISO 27001 Certification (Aspirational):",
    "   • Industry recognition",
    "   • Regulatory confidence",
    "   • Competitive differentiation",
    "",
    "🤖 Continuous AI Innovation:",
    "   • Expand MCP integrations",
    "   • Custom threat detection models"
])

# Slide 13: Partnership & Business Enablement
add_content_slide(prs, "Partnership & Business Enablement", [
    "Real Examples of Business Enablement:",
    "",
    "🚀 Market Launches:",
    "   • Colorado & Michigan: Regulatory compliance, security readiness",
    "   • Ontario (Q1 2026): Proactive security assessments in progress",
    "",
    "🏈 Super Bowl 2025 Scalability:",
    "   • Cribl + AKS scalability testing",
    "   • Validated 10-20x load handling",
    "   • Ensured platform resilience for peak events",
    "",
    "💰 Cost Optimization:",
    "   • 33% Sentinel cost reduction (freed budget for new tools)",
    "   • 40% WAF log storage savings",
    "",
    "How We Enable Speed:",
    "   ✓ Risk-based decisions (not approval/denial)",
    "   ✓ Automation reduces manual gates",
    "   ✓ Proactive security embedded early",
    "   ✓ 'We don't say no, we say here's the risk'"
])

# Slide 14: What We Need From You
add_content_slide(prs, "What We Need From You", [
    "🤝 Continued Partnership Across Teams:",
    "   • Security is everyone's responsibility",
    "   • Early engagement on new initiatives",
    "   • Transparent communication about risks",
    "",
    "💡 Embrace the FAIL Mindset:",
    "   • First Attempt In Learning",
    "   • Fail early, fail often, fail fast",
    "   • Innovation requires experimentation",
    "",
    "🎓 Security Awareness Participation:",
    "   • Complete NINJIO video training",
    "   • Report suspicious activity promptly",
    "   • Attend live security sessions",
    "",
    "💬 Feedback and Collaboration:",
    "   • Tell us what's working, what's not",
    "   • Help us improve processes",
    "   • Partner with us on solutions",
    "",
    "Together, we build a secure foundation for Hard Rock Digital's growth."
])

# Slide 15: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = VIOLET_1

# Q&A Title
qa_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
qa_frame = qa_box.text_frame
qa_frame.text = "Questions & Discussion"
qa_para = qa_frame.paragraphs[0]
qa_para.font.size = Pt(60)
qa_para.font.bold = True
qa_para.font.color.rgb = WHITE
qa_para.alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "Security Team | Hard Rock Digital"
contact_para = contact_frame.paragraphs[0]
contact_para.font.size = Pt(20)
contact_para.font.color.rgb = WHITE
contact_para.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "/home/user/test/State_of_Security_2026_Townhall.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Brand colors applied: Hard Rock Digital Violet 1 (#6a46f3)")
