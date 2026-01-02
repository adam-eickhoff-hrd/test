#!/usr/bin/env python3
"""
Create State of Security 2026 Townhall Presentation
Using the Hard Rock Digital PowerPoint Template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

# Load the template
template_path = "/home/user/test/HRD_PowerPoint-Template_v1.pptx"
prs = Presentation(template_path)

# Clear all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Helper function to set text in placeholder
def set_placeholder_text(slide, placeholder_idx, text, font_size=None, bold=False):
    """Set text in a placeholder by index"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            shape.text = text
            if font_size or bold:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if font_size:
                            run.font.size = Pt(font_size)
                        if bold:
                            run.font.bold = True
            return shape
    return None

def add_bullet_points(text_frame, items, font_size=18):
    """Add bullet points to a text frame"""
    text_frame.clear()
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0 if not item.startswith('   •') else 1
        p.font.size = Pt(font_size if p.level == 0 else font_size - 2)
        p.space_before = Pt(6)

# Slide 1: Title Slide - Layout [0] White-Title Plain
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_placeholder_text(slide, 0, "State of Security 2026", font_size=48, bold=True)
set_placeholder_text(slide, 1, "Building Tomorrow's Security\nHard Rock Digital Security Townhall", font_size=24)

# Slide 2: Agenda - Layout [1] White-Title Left
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Agenda", font_size=36, bold=True)
text_placeholder = None
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        text_placeholder = shape
        break
if text_placeholder:
    add_bullet_points(text_placeholder.text_frame, [
        "Year in Review: 2025 Achievements",
        "Deep Dive: AI-Powered Security Operations",
        "Security Operations & Detection Engineering",
        "Governance & Maturity",
        "Current State: Where We Stand Today",
        "Future State: Zero Trust Vision",
        "2026 Priorities & 3-Year Roadmap",
        "Our Philosophy: Partnership & Risk-Based Security",
        "Q&A"
    ])

# Slide 3: Our Security Philosophy
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Our Security Philosophy", font_size=36, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🤝 Risk-Based Approach",
            "   • We inform, business decides",
            "   • Security highlights risks, leadership accepts/mitigates",
            "",
            "🚀 Business Enablement",
            "   • Go as fast as you can possibly go",
            "   • We remove blockers, not create them",
            "",
            "💡 FAIL = First Attempt In Learning",
            "   • Fail early, fail often, fail fast",
            "   • Learning mindset over perfection",
            "",
            "🤝 Partnership, Not Gatekeeping",
            "   • We don't approve/deny—we collaborate"
        ], font_size=16)

# Slide 4: 2025 By The Numbers
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "2025 By The Numbers", font_size=36, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "📊 Key Metrics:",
            "   • 500 analytics rules deployed (at platform capacity)",
            "   • 79.7% alert closure rate maintained",
            "   • 49+ vendor security reviews completed",
            "   • ~560 employees trained on security awareness",
            "   • 33% cost reduction in Sentinel logging",
            "",
            "👥 Team Growth:",
            "   • Director of SRM hired",
            "   • 3 new SOC analysts, 2 new SRM analysts",
            "   • 1 Security Engineer (Jan 2026)",
            "",
            "🌎 Market Launches:",
            "   • Colorado & Michigan successfully launched",
            "   • Ontario launch prep underway (Q1 2026)"
        ], font_size=16)

# Slide 5: AI & MCP: Leading the Industry (Part 1)
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "🤖 AI & MCP: Industry Pioneers", font_size=32, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "What is Model Context Protocol (MCP)?",
            "   • Standardized way for AI to interact with external systems",
            "   • Real-time access to security tools and data",
            "   • Natural language interface for complex operations",
            "",
            "Why It Matters:",
            "   • We're among the FIRST security teams deploying production MCP",
            "   • Transforms how analysts interact with security tools",
            "   • Reduces manual work, accelerates investigation",
            "   • Enables AI-assisted decision making",
            "",
            "Our Position:",
            "   • Industry leadership in AI-powered security operations",
            "   • Innovation differentiator for Hard Rock Digital",
            "   • Competitive advantage in threat detection & response"
        ], font_size=15)

# Slide 6: AI & MCP: Our Implementations - Two Column Layout [8]
slide = prs.slides.add_slide(prs.slide_layouts[8])
set_placeholder_text(slide, 0, "🤖 AI & MCP: Our Implementations", font_size=32, bold=True)

# Left column
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "TheHive MCP Server:",
            "• Natural language case management",
            "• AI-assisted investigation workflow",
            "• Real-time alert context generation",
            "",
            "Vanta MCP Server:",
            "• Programmatic compliance access",
            "• Automated security posture reporting",
            "",
            "AI SOC Level 1 'Analyst':",
            "• Initial alert triage automation",
            "• Pattern recognition across alerts",
            "• Reduces analyst burnout"
        ], font_size=14)
        break

# Right column
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 2:
        add_bullet_points(shape.text_frame, [
            "Current AI Capabilities:",
            "• Claude AI integration for analysis",
            "• AI-powered log analysis",
            "• RSS cyber threat summarization",
            "",
            "2026 AI Roadmap:",
            "• Expand MCP to additional platforms",
            "• AI-powered playbook automation",
            "• Custom threat detection models",
            "",
            "Governance:",
            "• Hallucination testing completed",
            "• Least privilege OAuth credentials",
            "• Security evaluation framework"
        ], font_size=14)
        break

# Slide 7: Security Operations Excellence
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Security Operations Excellence", font_size=32, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🔒 24/7 Monitoring - Continent8 MSOC Partnership:",
            "   • 13 SOC analyst accounts onboarded",
            "   • API integration for automated alert forwarding",
            "   • n8n workflow automation (6pm-2am EST coverage)",
            "",
            "📈 Detection Engineering:",
            "   • 500 analytics rules at platform capacity",
            "   • 79.7% alert closure rate",
            "   • Major tuning: AWS false positives reduced from 647 to 4 events",
            "   • New capabilities: DPRK email detection, process hollowing",
            "",
            "🎯 Incident Response Wins:",
            "   • GraphQL scraping: Multiple attacks blocked, zero customer impact",
            "   • Credential stuffing (Dec 2025): 857 users rate-limited, no compromises",
            "   • Super Bowl 2025 prep: Validated 10-20x load handling"
        ], font_size=15)

# Slide 8: Governance & Maturity Milestones
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Governance & Maturity Milestones", font_size=32, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "📋 ISMS Policies Approved (December 2025):",
            "   • Security & Privacy Steering Committee re-established",
            "   • Top Level Security Policy, ISMS Roles & Responsibilities",
            "   • Patch Management SLAs, Backup Policy, Vulnerability Scanning",
            "",
            "🏆 ISO 22301 BCMS Framework:",
            "   • Comprehensive Business Continuity Management System",
            "   • 10 BCMS documents completed",
            "   • RTO/RPO targets established (Tier 1: <1hr, Tier 2: <4hrs, Tier 3: <24hrs)",
            "   • Gaming Operations, Payment Processing, Cybersecurity plans",
            "",
            "✅ ISO 27001 Gap Analysis:",
            "   • 96 controls mapped across 36 systems",
            "   • RACI matrix development",
            "   • Foundation for ISO 27001 certification (aspirational)"
        ], font_size=15)

# Slide 9: Current State - Two Column
slide = prs.slides.add_slide(prs.slide_layouts[8])
set_placeholder_text(slide, 0, "Current State: Where We Stand", font_size=32, bold=True)

# Left column - Strengths
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "✅ Strengths:",
            "• Governance maturity (ISMS approved)",
            "• AI/MCP industry leadership",
            "• 24/7 monitoring with MSOC",
            "• Strong team growth (6 new hires)",
            "• Cost optimization (33% savings)",
            "• Business enablement (CO, MI launches)",
            "• 79.7% alert closure rate",
            "",
            "🎯 Mature Capabilities:",
            "• Detection engineering",
            "• Incident response",
            "• Vendor risk management",
            "• Business continuity planning"
        ], font_size=13)
        break

# Right column - Opportunities
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 2:
        add_bullet_points(shape.text_frame, [
            "🎯 Opportunities (2026 Focus):",
            "• Sentinel capacity (at 500 rule limit)",
            "• Platform integrations:",
            "   - Teleskope DLP (Q1 2026)",
            "   - Sublime Security email (Q1 2026)",
            "   - Flare.io dark web (Q1 2026)",
            "• Identity governance (Entra ID rollout)",
            "• Penetration testing cadence",
            "• Ontario launch security readiness",
            "",
            "💡 Not Weaknesses—Growth Areas:",
            "• Every gap has a funded plan",
            "• Proactive identification shows maturity",
            "• Risk-based prioritization"
        ], font_size=13)
        break

# Slide 10: Future State Vision: Zero Trust
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Future State Vision: Zero Trust", font_size=32, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🎯 Zero Trust Principles:",
            "   • Never trust, always verify",
            "   • Identity-based access (not perimeter-based)",
            "   • Least privilege by default",
            "   • Continuous authentication & authorization",
            "",
            "✅ Current Progress:",
            "   • Cloudflare Zero Trust implementation (identity-based access)",
            "   • Entra Conditional Access policies (12-hour session limits)",
            "   • Admin account concept across teams",
            "   • Island Enterprise Browser (Phase 1 deployed)",
            "",
            "🚀 The Journey Ahead:",
            "   • Complete Island Browser enterprise rollout",
            "   • Expand identity-based controls",
            "   • Implement continuous authentication",
            "   • Enhance micro-segmentation"
        ], font_size=15)

# Slide 11: 2026 Priorities (Q1-Q2)
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "2026 Priorities: Q1-Q2", font_size=36, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🌎 Ontario Market Launch Security Readiness",
            "   • Complete security assessments & regulatory documentation",
            "   • Ensure compliance with Ontario gaming regulations",
            "",
            "🎯 ISO 22301 BCMS Execution",
            "   • Execute recovery strategies, conduct tabletop exercises",
            "   • Validate RTO/RPO targets through live testing",
            "   • Pursue ISO 22301 certification (competitive differentiator)",
            "",
            "📊 Sentinel Capacity Resolution",
            "   • Address 500 analytics rule limit blocking new detections",
            "",
            "🔒 Penetration Testing Program",
            "   • TrustedSec LLM assessment (January 20, 2026)",
            "   • Establish recurring penetration test cadence",
            "",
            "🚀 New Platform Integrations",
            "   • Teleskope DLP, Sublime Security email, Flare.io dark web"
        ], font_size=15)

# Slide 12: 3-Year Vision (2026-2028)
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "3-Year Vision: 2026-2028", font_size=36, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🔑 Identity Security Focus:",
            "   • Comprehensive identity governance (Entra ID)",
            "   • Quarterly User Access Reviews across 20+ systems",
            "   • Role-based access control (RBAC) maturity",
            "",
            "📧 Email Security Maturity:",
            "   • Sublime Security full deployment",
            "   • Advanced phishing protection",
            "   • Business email compromise prevention",
            "",
            "🌍 International Expansion Readiness:",
            "   • Adapt to shifting business priorities",
            "   • Security frameworks for new jurisdictions beyond Ontario",
            "",
            "🏆 ISO 27001 Certification (Aspirational):",
            "   • Industry recognition, Regulatory confidence",
            "",
            "🤖 Continuous AI Innovation:",
            "   • Expand MCP integrations, Custom threat detection models"
        ], font_size=15)

# Slide 13: Partnership & Business Enablement
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "Partnership & Business Enablement", font_size=30, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "Real Examples of Business Enablement:",
            "",
            "🚀 Market Launches:",
            "   • Colorado & Michigan: Regulatory compliance, security readiness",
            "   • Ontario (Q1 2026): Proactive security assessments in progress",
            "",
            "🏈 Super Bowl 2025 Scalability:",
            "   • Cribl + AKS scalability testing",
            "   • Validated 10-20x load handling",
            "",
            "💰 Cost Optimization:",
            "   • 33% Sentinel cost reduction (freed budget for new tools)",
            "",
            "How We Enable Speed:",
            "   ✓ Risk-based decisions (not approval/denial)",
            "   ✓ Automation reduces manual gates",
            "   ✓ 'We don't say no, we say here's the risk'"
        ], font_size=15)

# Slide 14: What We Need From You
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_placeholder_text(slide, 0, "What We Need From You", font_size=36, bold=True)
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        add_bullet_points(shape.text_frame, [
            "🤝 Continued Partnership Across Teams:",
            "   • Security is everyone's responsibility",
            "   • Early engagement on new initiatives",
            "",
            "💡 Embrace the FAIL Mindset:",
            "   • First Attempt In Learning",
            "   • Fail early, fail often, fail fast",
            "   • Innovation requires experimentation",
            "",
            "🎓 Security Awareness Participation:",
            "   • Complete NINJIO video training",
            "   • Report suspicious activity promptly",
            "",
            "💬 Feedback and Collaboration:",
            "   • Tell us what's working, what's not",
            "   • Help us improve processes",
            "   • Partner with us on solutions",
            "",
            "Together, we build a secure foundation for Hard Rock Digital's growth."
        ], font_size=15)

# Slide 15: Q&A - Use Violet Thank You layout [19]
slide = prs.slides.add_slide(prs.slide_layouts[19])

# Add large Q&A text
from pptx.enum.shapes import MSO_SHAPE
qa_shape = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
text_frame = qa_shape.text_frame
text_frame.text = "Questions & Discussion"
p = text_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "/home/user/test/State_of_Security_2026_Townhall_v2.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Using Hard Rock Digital official template layouts")
print(f"📐 Slide dimensions: 13.33\" x 7.50\" (16:9)")
